#!/usr/bin/env python3

import csv
import os
from collections import Counter
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTDIR = os.path.join("data", "processed", "analysis", "2026", "R1")
PLOTS_DIR = os.path.join("plots", "2026", "R1")
CSV_PATH = os.path.join("data", "processed", "transcriptions", "output", "2026", "transcripcion_R1.csv")
PPT_PATH = os.path.join(OUTDIR, "R1_2026_composicion_panel.pptx")


def read_basic_stats(csv_path: str):
    rows = 0
    segmentos = Counter()
    with open(csv_path, "r", encoding="utf-8-sig", newline="") as f:
        reader = csv.DictReader(f)
        for row in reader:
            rows += 1
            seg = (row.get("segmento") or "").strip() or "Sin dato"
            segmentos[seg] += 1
    return rows, segmentos


def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False,
               color=(34, 34, 34), align=PP_ALIGN.LEFT, font_name="Aptos"):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = RGBColor(*color)
    return tx


def add_cover(prs, n_total, segmentos_counter):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_shape(
        autoshape_type_id=1,  # rectangle
        left=Inches(0),
        top=Inches(0),
        width=prs.slide_width,
        height=Inches(1.15),
    ).fill.solid()
    banner = slide.shapes[-1]
    banner.fill.fore_color.rgb = RGBColor(20, 44, 78)
    banner.line.fill.background()

    add_textbox(
        slide, Inches(0.6), Inches(0.3), Inches(11.9), Inches(0.6),
        "R1 2026 - Composicion del panel",
        font_size=24, bold=True, color=(255, 255, 255)
    )

    add_textbox(
        slide, Inches(0.7), Inches(1.45), Inches(12.0), Inches(0.8),
        "Resumen descriptivo del panel a partir de transcripcion_R1.csv",
        font_size=19, color=(60, 60, 60)
    )

    items = [
        f"N total de casos: {n_total}",
        f"Segmentos presentes: {len(segmentos_counter)}",
        "Graficos generados con ggplot2",
        f"Fecha de generacion: {date.today().isoformat()}",
    ]

    y = 2.35
    for item in items:
        add_textbox(
            slide, Inches(1.0), Inches(y), Inches(11.0), Inches(0.45),
            f"- {item}", font_size=18
        )
        y += 0.5

    add_textbox(
        slide, Inches(0.8), Inches(5.0), Inches(5.7), Inches(1.7),
        "Contenido\n1. Segmentos\n2. Genero\n3. Edad\n4. Educacion\n5. Variables complementarias",
        font_size=16, bold=False, color=(30, 30, 30), font_name="Aptos"
    )

    top_seg = segmentos_counter.most_common()
    seg_lines = [f"{name}: {count}" for name, count in top_seg]
    add_textbox(
        slide, Inches(6.3), Inches(5.0), Inches(5.9), Inches(1.9),
        "Segmentos (n)\n" + "\n".join(seg_lines),
        font_size=15, color=(30, 30, 30)
    )

    add_textbox(
        slide, Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.25),
        f"Fuente: {CSV_PATH}",
        font_size=9, color=(100, 100, 100)
    )


def add_image_slide(prs, title, image_path, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.45),
               title, font_size=22, bold=True, color=(20, 44, 78))
    if subtitle:
        add_textbox(slide, Inches(0.45), Inches(0.65), Inches(12.4), Inches(0.3),
                   subtitle, font_size=11, color=(90, 90, 90))

    slide.shapes.add_picture(image_path, Inches(0.35), Inches(0.95), width=Inches(12.65))
    add_textbox(slide, Inches(0.45), Inches(7.05), Inches(12.2), Inches(0.2),
               f"Fuente: {CSV_PATH}", font_size=9, color=(105, 105, 105))


def main():
    n_total, segmentos = read_basic_stats(CSV_PATH)

    required_images = [
        "01_segmento.png",
        "02_genero.png",
        "03_edad.png",
        "04_educacion.png",
        "05_complementarias.png",
    ]
    missing = [img for img in required_images if not os.path.exists(os.path.join(PLOTS_DIR, img))]
    if missing:
        raise FileNotFoundError(f"Faltan graficos: {missing}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs, n_total, segmentos)
    add_image_slide(prs, "1. Composicion por segmento", os.path.join(PLOTS_DIR, "01_segmento.png"))
    add_image_slide(prs, "2. Composicion por genero", os.path.join(PLOTS_DIR, "02_genero.png"))
    add_image_slide(prs, "3. Composicion por edad", os.path.join(PLOTS_DIR, "03_edad.png"))
    add_image_slide(prs, "4. Composicion por nivel educativo", os.path.join(PLOTS_DIR, "04_educacion.png"))
    add_image_slide(prs, "5. Variables complementarias", os.path.join(PLOTS_DIR, "05_complementarias.png"),
                    subtitle="Departamento (top 10 + otros), voto y etiqueta")

    os.makedirs(OUTDIR, exist_ok=True)
    prs.save(PPT_PATH)
    print(PPT_PATH)


if __name__ == "__main__":
    main()
